from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageFont, ImageDraw
import make_quizz
import pts
from dotenv import load_dotenv
from io import BytesIO
from supabase import create_client, Client
import os
from datetime import datetime
import numpy as np
import pandas as pd
from wordcloud import WordCloud, STOPWORDS, ImageColorGenerator
import matplotlib.pyplot as plt
import warnings
from IPython.display import display
import pptx
import reference


load_dotenv()
url: str = os.environ.get("SUPABASE_URL")
key: str = os.environ.get("SUPABASE_KEY")
supabase: Client = create_client(url, key)

class PPT:
    def __init__(self):
        self.presentation = Presentation("template.pptx")

    def add_title_slide(self, title_data):
        slide = self.presentation.slides[0]
        for shape, data in zip(slide.shapes, [title_data]):
            if not shape.has_text_frame:
                continue
            shape.text_frame.text = data
            for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(38, 38, 38)
                        run.font.size = Pt(26)
                        run.font.bold = True

    def add_content_slide(self, title_data, content_data, i):
        slide = self.presentation.slides[i]
        cnt = 0
        for shape, data in zip(slide.shapes, [title_data,content_data]):
            if not shape.has_text_frame:
                continue
            shape.text_frame.text = data
            for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(26)
                        if cnt==0:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.bold = True
                            cnt+=1
                        else:
                            run.font.color.rgb = RGBColor(0, 0, 0)
                            run.font.bold = False

    def add_image_slide(self, title_data, img_path, i):
        slide = self.presentation.slides[i]
        cnt = 0
        pic = slide.shapes.add_picture(img_path, pptx.util.Inches(2), pptx.util.Inches(2),width=pptx.util.Inches(9), height=pptx.util.Inches(5))


    def create_wordcloud(self, summary_string, name):
        wordcloud = WordCloud().generate(summary_string)

        plt.figure()
        plt.imshow(wordcloud, interpolation="bilinear")
        plt.axis("off")

        # Save the plot to a file
        plt.savefig(name, bbox_inches='tight', pad_inches=0.2)

        # Optional: Close the plot to release resources (comment this line if you want to display the plot)
        plt.close()

    def create_presentation(self, path):
        summary_string = pts.PdfToString(path)
        wordcloud_image_name = datetime.now().strftime("%Y%m%d%H%M%S%f") + f"{summary_string[:3]}.png"
        self.create_wordcloud(summary_string, wordcloud_image_name)

        supabase.storage.from_("url").upload(file=wordcloud_image_name, path=wordcloud_image_name, file_options={"content-type": "image/png"})
        res_image = supabase.storage.from_('url').get_public_url(wordcloud_image_name)


        """논문 작성에 참조된 논문 리스트 뽑아오기"""
        reference_list = summary_string.split("\n")
        reference_data = reference.get_reference(reference_list)
        
        summary_string = make_quizz.summary(summary_string)
        summary_string_list = summary_string.split("\n\n")

        title_name = ""
        for i in range(len(summary_string_list)):
            if i == 0:
                self.add_title_slide(summary_string_list[i].split(": ")[1])
                img = Image.open("background.png")
                img_width, img_height = img.size
                draw = ImageDraw.Draw(img)
                font_size = 30
                font = ImageFont.truetype("NanumGothic.ttf", font_size)
                text_width = draw.textlength(summary_string_list[i].split(": ")[1], font)
                x = (img_width-text_width)//2
                y = (img_height-font_size)//2
                draw.text((x,y),summary_string_list[i].split(": ")[1], (0,0,0), font=font)
                draw = ImageDraw.Draw(img)
                title_name = datetime.now().strftime("%Y%m%d%H%M%S%f") + ".png"
                img.save(title_name)
            elif i > 0 and i < len(summary_string_list):
                data = summary_string_list[i].split(": ")
                self.add_content_slide(data[0], data[1], i)
        
        self.add_image_slide("wordcloud",wordcloud_image_name,9)
        self.add_content_slide("References", '\n'.join(reference_data[1]), 10)

        supabase.storage.from_("url").upload(file=title_name, path=title_name, file_options={"content-type": "image/png"})
        title_image = supabase.storage.from_('url').get_public_url(title_name)

        timestamp1 = datetime.now().strftime("%Y%m%d%H%M%S%f")
        name = f"{timestamp1}.pptx"
        self.presentation.save(name)

        supabase.storage.from_("url").upload(file=name, path=name, file_options={"content-type": "application/vnd.ms-powerpoint"})
        res = supabase.storage.from_('url').get_public_url(name)

        return {"res": res, "res_image": res_image, "title_image": title_image}

# Example usage:

#ppt_generator = PPT()
#ppt_generator.create_presentation("ue.pdf")
